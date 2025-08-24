import os
import zipfile
import shutil
from pptx import Presentation
from pathlib import Path
import patoolib
import argparse


class PPTXProcessor:
    """
    处理PPTX文件和压缩文件的工具类

    功能：
    1. 递归解压目录中的所有ZIP和RAR文件（可选择删除原始压缩文件）
    2. 删除所有PPTX文件的最后一页
    """

    def __init__(self, root_dir, delete_archive=False):
        """
        初始化处理器

        :param root_dir: 要处理的根目录
        :param delete_archive: 解压后是否删除原始压缩文件
        """
        self.root_dir = Path(root_dir).resolve()
        self.delete_archive = delete_archive
        if not self.root_dir.exists():
            raise FileNotFoundError(f"目录不存在: {self.root_dir}")

    def process_all(self):
        """
        执行所有处理步骤
        """
        print(f"开始处理目录: {self.root_dir}")
        self.extract_all_archives()
        self.process_all_pptx()
        print("处理完成!")

    def extract_all_archives(self):
        """
        递归解压所有ZIP和RAR文件
        """
        for root, _, files in os.walk(self.root_dir):
            for file in files:
                file_path = Path(root) / file
                if file.lower().endswith('.zip'):
                    self._extract_zip(file_path)
                elif file.lower().endswith('.rar'):
                    self._extract_rar(file_path)

    def _extract_zip(self, zip_path):
        """
        解压ZIP文件到当前目录

        :param zip_path: ZIP文件路径
        """
        print(f"解压ZIP文件: {zip_path}")
        extract_dir = zip_path.parent

        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                # 先检查是否有同名文件冲突
                conflict_files = self._check_conflicts(zip_ref, extract_dir)
                if conflict_files:
                    print(f"警告: 解压将覆盖以下已存在的文件: {', '.join(conflict_files)}")

                zip_ref.extractall(extract_dir)
            print(f"解压完成到: {extract_dir}")

            if self.delete_archive:
                self._safe_delete(zip_path)

        except Exception as e:
            print(f"解压ZIP文件失败 {zip_path}: {str(e)}")

    def _extract_rar(self, rar_path):
        """
        解压RAR文件到当前目录

        :param rar_path: RAR文件路径
        """
        print(f"解压RAR文件: {rar_path}")
        extract_dir = rar_path.parent

        try:
            # patoolib没有提供直接检查冲突的方法，所以先解压再处理
            temp_dir = extract_dir / f"temp_{rar_path.stem}"
            patoolib.extract_archive(str(rar_path), outdir=str(temp_dir))

            # 移动文件并处理冲突
            self._move_contents(temp_dir, extract_dir)

            # 删除临时目录
            shutil.rmtree(temp_dir)
            print(f"解压完成到: {extract_dir}")

            if self.delete_archive:
                self._safe_delete(rar_path)

        except Exception as e:
            print(f"解压RAR文件失败 {rar_path}: {str(e)}")
            if temp_dir.exists():
                shutil.rmtree(temp_dir, ignore_errors=True)

    def _check_conflicts(self, zip_ref, target_dir):
        """
        检查解压是否会覆盖现有文件
        :param zip_ref: ZipFile对象
        :param target_dir: 目标目录
        :return: 冲突文件列表
        """
        conflict_files = []
        for file in zip_ref.namelist():
            if not file.endswith('/'):  # 忽略目录
                dest_path = Path(target_dir) / file
                if dest_path.exists():
                    conflict_files.append(file)
        return conflict_files

    def _move_contents(self, src_dir, dest_dir):
        """
        移动源目录所有内容到目标目录，处理冲突
        :param src_dir: 源目录
        :param dest_dir: 目标目录
        """
        for item in src_dir.iterdir():
            dest_item = dest_dir / item.name
            if dest_item.exists():
                if dest_item.is_dir():
                    # 如果是目录，递归合并
                    self._move_contents(item, dest_item)
                else:
                    # 如果是文件，覆盖
                    shutil.move(str(item), str(dest_item))
            else:
                shutil.move(str(item), str(dest_dir))

    def _safe_delete(self, file_path):
        """
        安全删除文件
        :param file_path: 要删除的文件路径
        """
        try:
            file_path.unlink()
            print(f"已删除原始压缩文件: {file_path}")
        except Exception as e:
            print(f"删除文件失败 {file_path}: {str(e)}")

    def process_all_pptx(self):
        """
        递归处理所有PPTX文件，删除最后一页
        """
        for root, _, files in os.walk(self.root_dir):
            for file in files:
                if file.lower().endswith('.pptx'):
                    pptx_path = Path(root) / file
                    self._remove_last_slide(pptx_path)

    def _remove_last_slide(self, pptx_path):
        """
        删除PPTX文件的最后一页

        :param pptx_path: PPTX文件路径
        """
        print(f"处理PPTX文件: {pptx_path}")
        backup_path = None

        try:
            # 创建备份文件
            backup_path = pptx_path.with_suffix('.pptx.bak')
            shutil.copy2(pptx_path, backup_path)

            # 加载演示文稿
            prs = Presentation(pptx_path)

            if len(prs.slides) == 0:
                print(f"警告: {pptx_path} 中没有幻灯片")
                return

            # 删除最后一页
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            last_slide = slides[-1]
            xml_slides.remove(last_slide)

            # 保存修改后的文件
            prs.save(pptx_path)
            print(f"已删除最后一页并保存: {pptx_path}")

        except Exception as e:
            print(f"处理PPTX文件失败 {pptx_path}: {str(e)}")
            # 如果出错，恢复备份
            if backup_path and backup_path.exists():
                shutil.move(backup_path, pptx_path)
        finally:
            # 如果处理成功，删除备份
            if backup_path and backup_path.exists():
                try:
                    backup_path.unlink()
                except:
                    pass


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX处理工具")
    parser.add_argument("directory", help="要处理的目录路径")
    parser.add_argument("--delete", action="store_true",
                        help="解压后删除原始压缩文件")
    args = parser.parse_args()

    processor = PPTXProcessor(args.directory, args.delete)
    processor.process_all()


# python .\unzip-files.py D:\workspace\资源\test
